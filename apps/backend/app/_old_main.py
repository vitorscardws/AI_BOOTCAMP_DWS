
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer, util
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
from langdetect import detect
from pptx import Presentation
import os
import fitz  # PyMuPDF
import pickle
import requests

load_dotenv()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# PDF_PATH = "docs/MTG_BASIC_RULES.pdf"
# CACHE_PATH = "docs/mtg_embeddings.pkl"
PDF_PATH = "docs/MAGIC_RULES.pdf"
PPTX_PATH = "docs/PRESENTATION.pptx"
CACHE_PATH = "docs/source.pkl"
CHUNK_SIZE = 500

CACHE_PATH_PDF = "docs/pdf_embeddings.pkl"
CACHE_PATH_PPTX = "docs/pptx_embeddings.pkl"

# OpenAI API configuration
api_key = os.getenv("OPENAI_API_KEY")
url = "https://api.openai.com/v1/chat/completions"
headers = {
    "Authorization": f"Bearer {api_key}",
    "Content-Type": "application/json"
}
openai_model = "gpt-4"  # or "gpt-3.5-turbo" if preferred

# GROQ API configuration
api_key_groq = os.getenv("GROQ_API_KEY")
groq_url = "https://api.groq.com/openai/v1/chat/completions"
groq_headers = {
    "Authorization": f"Bearer {api_key_groq}",
    "Content-Type": "application/json"
}
groq_model = "llama-3.3-70b-versatile"  



# Model
model = SentenceTransformer("all-MiniLM-L6-v2")

# Request schema
class QueryRequest(BaseModel):
    query: str

# Utils
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    return "\n".join([page.get_text() for page in doc])

def ppt_to_text(file_path):
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text)

def chunk_text(text, chunk_size):
    chunks, current = [], ""
    for paragraph in text.split("\n"):
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        if len(current) + len(paragraph) <= chunk_size:
            current += " " + paragraph
        else:
            chunks.append(current.strip())
            current = paragraph
    if current:
        chunks.append(current.strip())
    return chunks

def load_or_create_embeddings_pdf():
    if os.path.exists(CACHE_PATH_PDF):
        with open(CACHE_PATH_PDF, "rb") as f:
            return pickle.load(f)

    if not os.path.exists(PDF_PATH):
        raise FileNotFoundError(f"Arquivo PDF não encontrado: {PDF_PATH}")

    text = extract_text_from_pdf(PDF_PATH)
    chunks = chunk_text(text, CHUNK_SIZE)
    documents = [{"id": f"pdf_{i}", "text": chunk} for i, chunk in enumerate(chunks)]
    doc_embeddings = {
        doc["id"]: model.encode(doc["text"], convert_to_tensor=True)
        for doc in documents
    }

    with open(CACHE_PATH_PDF, "wb") as f:
        pickle.dump((documents, doc_embeddings), f)

    return documents, doc_embeddings


def load_or_create_embeddings_pptx():
    if os.path.exists(CACHE_PATH_PPTX):
        with open(CACHE_PATH_PPTX, "rb") as f:
            return pickle.load(f)

    if not os.path.exists(PPTX_PATH):
        raise FileNotFoundError(f"Arquivo PPTX não encontrado: {PPTX_PATH}")

    text = ppt_to_text(PPTX_PATH)
    chunks = chunk_text(text, CHUNK_SIZE)
    documents = [{"id": f"ppt_{i}", "text": chunk} for i, chunk in enumerate(chunks)]
    doc_embeddings = {
        doc["id"]: model.encode(doc["text"], convert_to_tensor=True)
        for doc in documents
    }

    with open(CACHE_PATH_PPTX, "wb") as f:
        pickle.dump((documents, doc_embeddings), f)

    return documents, doc_embeddings



# Pré-carrega embeddings na inicialização
pdf_docs, pdf_embeddings = load_or_create_embeddings_pdf()
ppt_docs, ppt_embeddings = load_or_create_embeddings_pptx()



# Endpoint
@app.post("/query")
async def query(request: QueryRequest):
    query_embedding = model.encode(request.query, convert_to_tensor=True)

    best_doc = {}
    best_score = float('-inf')

    for doc in pdf_docs:
        score = util.cos_sim(query_embedding, pdf_embeddings[doc["id"]]).item()
        if score > best_score:
            best_score = score
            best_doc = doc

    if not best_doc:
        raise HTTPException(status_code=404, detail="No relevant document found.")
    

    print(f"Best document: {best_doc['text']}")  # Debugging line to see the best document

    language = detect(request.query)


    messages = [
        {
            "role": "system",
            "content": (
                "You are a Magic: The Gathering expert helping new players. "
                "Answer using only the provided content, but never mention or hint at the source or the text. "
                "Preserve the language of the user's question (if the question is in Portuguese, respond in Portuguese). "
                "Use a clear, friendly, and conversational tone. "
                "Be concise and objective, focusing only on what was asked. "
                "Do not guess or add information not found in the content."
            )
        },
        {
            "role": "user",
            "content": f"""
            {best_doc['text']}

            Question:
            {request.query}
            """
        }
    ]


    try:
        data = {
            "model": groq_model,
            "messages": messages,
            "temperature": 0.2,     
        }

        response = requests.post(groq_url, headers=groq_headers, json=data)
        res = response.json()
       
        
        return {"response": res.get("choices")[0].get("message").get("content")}


    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/text-to-mongo")
async def text_to_mongo(request: QueryRequest):


    messages = [
        {
            "role": "system",
            "content": (
                "You are a MongoDB expert. Your task is to convert natural language requests into valid MongoDB queries in JSON format.\n"
                "You MUST follow these strict rules:\n"
                "1. Only return the query, no explanations or markdown.\n"
                "2. You can interpret common business, financial, or database-related terminology, including date filters and numeric comparisons.\n"
                "3. NEVER assume field names — only use field names that are directly mentioned or clearly implied by the request.\n"
                "4. If you cannot confidently generate a valid query, respond with exactly:\n"
                '"Unable to generate a valid MongoDB query from the input."'
            )
        },
        {
            "role": "user",
            "content": (
                f'Convert the following request into a valid MongoDB query. '
                f'Input: "{request}"\n\n'
                "Answer:"
            )
        }
    ]
#note: use temperature 0.0 to get a more deterministic response
    try:
        data = {
             "model":  "gpt-3.5-turbo", # groq_model,
            "messages": messages,
            "temperature": 0.0
        }
        #response = requests.post(groq_url, headers=groq_headers, json=data)
        response = requests.post(url, headers=headers, json=data)
        result = response.json()

        mongo_query = result['choices'][0]['message']['content'].strip()
        return {"response": mongo_query}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    

@app.post("/ppt-search")
async def ppt_search(request: QueryRequest):

    query_embedding = model.encode(request.query, convert_to_tensor=True)

    best_doc = {}
    best_score = float('-inf')

    for doc in ppt_docs:
        score = util.cos_sim(query_embedding, ppt_embeddings[doc["id"]]).item()
        if score > best_score:
            best_score = score
            best_doc = doc

    if not best_doc:
        raise HTTPException(status_code=404, detail="No relevant document found.")
    

    print(f"Best document: {best_doc['text']}")

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert communicator who summarizes and explains the content of a PowerPoint presentation "
                "in a clear, engaging, and easy-to-understand way. "
                "Use only the information provided in the document text — do not guess or add information that is not present. "
                "Rewrite bullet points or fragmented text into complete, well-structured sentences. "
                "**Preserve the language of the user's question (if the question is in Portuguese, respond in Portuguese)**. "
                "Do not mention slides, bullet points, or that this came from a presentation. "
                "If the answer is not found in the text, say so clearly."
            )
        },
        {
            "role": "user",
            "content": f"""
            {best_doc['text']}

            Question:
            {request.query}
            """
        }
    ]

    try:
        data = {
             "model":  "gpt-4o-mini",
            "messages": messages,
            "temperature": 0.4
        }
        #response = requests.post(groq_url, headers=groq_headers, json=data)
        response = requests.post(url, headers=headers, json=data)
        result = response.json()

        mongo_query = result['choices'][0]['message']['content'].strip()
        return {"response": mongo_query}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))