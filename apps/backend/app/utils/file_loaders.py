import fitz
from pptx import Presentation

def extract_text_from_pdf(pdf_path: str) -> str:
    doc = fitz.open(pdf_path)
    return "\n".join([page.get_text() for page in doc])

def ppt_to_text(file_path: str) -> str:
    prs = Presentation(file_path)
    text_list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_list.append(shape.text)
    return "\n".join(text_list)
